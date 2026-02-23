# -*- coding: utf-8 -*-
"""
AdStrategy AI 발표자료 생성기
==============================
1분 브리핑 + 라이브 데모 구조의 PPTX 슬라이드 자동 생성.

실행: python scripts/generate_presentation.py
출력: AdStrategy_AI_발표자료.pptx
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = os.path.join(ROOT, "figures")

# ── Design Tokens ────────────────────────────────────────────────────
C_DARK = RGBColor(0x0D, 0x11, 0x17)
C_CARD = RGBColor(0x16, 0x1B, 0x22)
C_LIGHT = RGBColor(0xF0, 0xF2, 0xF5)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_RED = RGBColor(0xE7, 0x4C, 0x3C)
C_BLUE = RGBColor(0x34, 0x98, 0xDB)
C_GREEN = RGBColor(0x2E, 0xCC, 0x71)
C_GRAY = RGBColor(0x8B, 0x94, 0x9E)
C_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
C_NAVY = RGBColor(0x0F, 0x34, 0x60)
C_DIM = RGBColor(0x60, 0x60, 0x60)

FONT = "Malgun Gothic"
FONT_CODE = "Consolas"
SW = Inches(13.333)
SH = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────
def _slide(prs, bg=C_DARK):
    try:
        layout = prs.slide_layouts[6]
    except IndexError:
        layout = prs.slide_layouts[-1]
    s = prs.slides.add_slide(layout)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def _accent_bar(s, color=C_RED):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(s, txt, x, y, w, h, sz=24, bold=False, color=C_WHITE,
          align=PP_ALIGN.LEFT, font=FONT):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf


def _multi(s, lines, x, y, w, h, color=C_WHITE, font=FONT,
           align=PP_ALIGN.LEFT):
    """lines: [(text, font_size, is_bold), ...]"""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(4)
    return tf


def _notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def _rect(s, x, y, w, h, fill, border=None, bw=Pt(0)):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = bw
    else:
        shape.line.fill.background()
    return shape


def _fig(s, filename, y=Inches(1.2), max_h=Inches(5.2), max_w=Inches(11)):
    path = os.path.join(FIGS, filename)
    if not os.path.exists(path):
        _text(s, f"[ {filename} — 파일 없음 ]",
              Inches(3), y + Inches(1.5), Inches(7), Inches(1),
              sz=18, color=C_GRAY, align=PP_ALIGN.CENTER)
        return None
    pic = s.shapes.add_picture(path, 0, y)
    iw, ih = pic.width, pic.height
    aspect = iw / ih
    w = int(max_w)
    h = int(w / aspect)
    if h > int(max_h):
        h = int(max_h)
        w = int(h * aspect)
    pic.width, pic.height = w, h
    pic.left = (int(SW) - w) // 2
    pic.top = y
    return pic


# ── Slide 1: Title ──────────────────────────────────────────────────
def s01_title(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)
    _text(s, "AdStrategy AI",
          Inches(1), Inches(2), Inches(11.3), Inches(1.2),
          sz=56, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _text(s, "데이터 감사로 찾은 정직한 성능",
          Inches(1), Inches(3.4), Inches(11.3), Inches(0.8),
          sz=28, color=C_GRAY, align=PP_ALIGN.CENTER)
    _text(s, "AI가 자기 모델의 거짓말을 잡아낸 이야기",
          Inches(1), Inches(5.0), Inches(11.3), Inches(0.6),
          sz=20, color=C_BLUE, align=PP_ALIGN.CENTER)
    _notes(s, "[타이틀] 3초 노출 후 다음 슬라이드로 전환.")


# ── Slide 2: Vision ──────────────────────────────────────────────────
def s02_vision(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)
    _text(s, "디지털 광고, AI가 설계해준다면?",
          Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
          sz=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    points = [
        ("$600B+", "글로벌 디지털 광고 시장", C_BLUE),
        ("3개 플랫폼", "Google · Meta · TikTok\n어디에 얼마를 써야 할까?", C_ORANGE),
        ("데이터 → 전략", "ML 예측 + AI 에이전트로\n광고 의사결정 자동화", C_GREEN),
    ]
    cw = Inches(3.4)
    xs = [Inches(1.3), Inches(5.0), Inches(8.7)]
    ty = Inches(1.8)
    for i, (big, desc, clr) in enumerate(points):
        _rect(s, xs[i], ty, cw, Inches(3.2), C_CARD, border=clr, bw=Pt(2))
        _text(s, big, xs[i], ty + Inches(0.4), cw, Inches(0.8),
              sz=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
        _text(s, desc, xs[i] + Inches(0.2), ty + Inches(1.5),
              cw - Inches(0.4), Inches(1.2),
              sz=15, color=C_GRAY, align=PP_ALIGN.CENTER)

    _text(s, "중소 광고주도 데이터 기반으로 전략을 세울 수 있는 도구를 만들자",
          Inches(1), Inches(5.6), Inches(11.3), Inches(0.6),
          sz=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _notes(s,
        "[나레이션]\n"
        "디지털 광고 시장은 600조 원이 넘는 거대한 시장입니다.\n"
        "Google, Meta, TikTok — 어디에, 얼마를, 언제 써야 하는지\n"
        "데이터 없이 결정하기 어렵습니다.\n"
        "그래서 ML 예측과 AI 에이전트로 이 의사결정을 자동화하는 도구를 만들기로 했습니다.")


# ── Slide 3: What We Built ──────────────────────────────────────────
def s03_built(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)
    _text(s, "무엇을 만들었나",
          Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
          sz=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    items = [
        ("데이터 파이프라인", "Kaggle 1,800건 →\n4단계 보강 → 10,030건\nFRED · World Bank · Trends",
         C_BLUE),
        ("ML 예측 모델", "ROAS · CPC · CPA 예측\nXGBoost + LightGBM 앙상블\nSHAP 설명가능성",
         C_RED),
        ("AI 에이전트", "GPT-4o Function Calling\n자연어 대화로 전략 설계\n4개 Tool 자동 호출",
         C_GREEN),
        ("Streamlit 웹 앱", "5탭 대시보드\n예산 시뮬레이터\n내 데이터 분석",
         C_ORANGE),
    ]
    box_w, box_h = Inches(2.5), Inches(3.0)
    xs = [Inches(0.6), Inches(3.5), Inches(6.4), Inches(9.3)]
    ty = Inches(1.6)
    for i, (title, desc, clr) in enumerate(items):
        _rect(s, xs[i], ty, box_w, box_h, C_CARD, border=clr, bw=Pt(2))
        _text(s, title, xs[i], ty + Inches(0.3), box_w, Inches(0.6),
              sz=16, bold=True, color=clr, align=PP_ALIGN.CENTER)
        _text(s, desc, xs[i] + Inches(0.15), ty + Inches(1.1),
              box_w - Inches(0.3), Inches(1.6),
              sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    _text(s, "→ 여기까지는 순조로웠습니다. 그런데...",
          Inches(1), Inches(5.4), Inches(11.3), Inches(0.6),
          sz=22, bold=True, color=C_RED, align=PP_ALIGN.CENTER)

    _notes(s,
        "[나레이션]\n"
        "만든 건 크게 네 가지입니다.\n"
        "첫째, Kaggle 원본 1,800건을 공공 API와 합성 데이터로 만 건까지 보강한 파이프라인.\n"
        "둘째, ROAS, CPC, CPA를 예측하는 앙상블 ML 모델.\n"
        "셋째, GPT-4o 기반으로 자연어 대화만으로 광고 전략을 설계해주는 AI 에이전트.\n"
        "넷째, 이 모든 걸 담은 Streamlit 웹 앱.\n\n"
        "여기까지는 순조로웠습니다. 그런데...")


# ── Slide 4: Hook ────────────────────────────────────────────────────
def s04_hook(prs):
    """원래 s02 — 번호만 변경"""
    s = _slide(prs, C_DARK)
    _accent_bar(s)
    _text(s, "R\u00b2 = 0.79",
          Inches(1), Inches(1.2), Inches(11.3), Inches(2),
          sz=84, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    _text(s, "\ub9c8\ucf00\ud305 \ub370\uc774\ud130\uc5d0\uc11c \uc774\uac74 \ub108\ubb34 \uc88b\uc740 \uc22b\uc790",
          Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
          sz=32, color=C_WHITE, align=PP_ALIGN.CENTER)
    _text(s, "\u201c\ub124\uac00 \uac70\uc9d3\ub9d0\ud558\uace0 \uc788\ub294 \uac74 \uc544\ub2c8\uc57c?\u201d",
          Inches(1), Inches(5.0), Inches(11.3), Inches(0.8),
          sz=28, color=C_GRAY, align=PP_ALIGN.CENTER)
    _notes(s,
        "[나레이션 0:00-0:10]\n"
        "AI에게 광고 성과를 예측하라고 했더니, R² 0.79가 나왔습니다.\n"
        "마케팅 데이터에서 이건 너무 좋은 숫자예요.\n"
        "그래서 이 AI에게 다시 물었습니다 — '네가 거짓말하고 있는 건 아니야?'")


# ── Slide 3: Pipeline ────────────────────────────────────────────────
def s03_pipeline(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)
    _text(s, "4단계 데이터 보강 파이프라인",
          Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
          sz=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    phases = [
        ("Phase 1", "Kaggle 병합\n1,800 \u2192 5,400건", C_BLUE),
        ("Phase 2", "외부 데이터\n거시경제\u00b7시즌\u00b7트렌드", C_BLUE),
        ("Phase 3", "크리에이티브\n메타데이터", C_RED),
        ("Phase 4", "경쟁+합성\n\u2192 10,030건", C_BLUE),
    ]
    bw, bh = Inches(2.5), Inches(2.0)
    xs = [Inches(0.6), Inches(3.5), Inches(6.4), Inches(9.3)]
    ty = Inches(1.8)

    for i, (title, desc, clr) in enumerate(phases):
        _rect(s, xs[i], ty, bw, bh, clr)
        _text(s, title, xs[i], ty + Inches(0.25), bw, Inches(0.5),
              sz=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _text(s, desc, xs[i], ty + Inches(0.85), bw, Inches(1.0),
              sz=15, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            ax = xs[i] + bw + Inches(0.05)
            _text(s, "\u2192", ax, ty + Inches(0.6), Inches(0.4), Inches(0.5),
                  sz=36, color=C_GRAY, align=PP_ALIGN.CENTER)

    _text(s, "\u26a0 Phase 3\uc5d0\uc11c \uc815\ub2f5\uc774 \ud53c\uccd0\uc5d0 \uc0c8\uc5b4 \ub4e4\uc5b4\uac10 (Target Leakage)",
          Inches(1), Inches(4.4), Inches(11.3), Inches(0.6),
          sz=22, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    _text(s, "bounce_rate, landing_page_load_time, creative_impact_factor \u2192 ROAS\uc758 \uc5ed\ud568\uc218",
          Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
          sz=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    _notes(s,
        "[나레이션 0:10-0:25]\n"
        "Kaggle 1,800건을 4단계 파이프라인으로 만 건까지 보강했는데,\n"
        "Phase 3에서 bounce_rate = 65 − ROAS × 2라는 공식으로\n"
        "정답이 피처에 새어 들어간 걸 찾았습니다.\n"
        "이른바 target leakage — 정답지를 미리 본 셈이죠.")


# ── Slide 4: Smoking Gun ─────────────────────────────────────────────
def s04_smoking_gun(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)
    _text(s, "The Smoking Gun",
          Inches(1), Inches(0.4), Inches(11.3), Inches(0.8),
          sz=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _rect(s, Inches(1.5), Inches(1.6), Inches(10.3), Inches(2.8),
          C_CARD, border=C_RED, bw=Pt(3))

    _text(s, "bounce_rate = 65 - ROAS \u00d7 2",
          Inches(2), Inches(1.9), Inches(9.3), Inches(0.8),
          sz=38, bold=True, color=C_RED, align=PP_ALIGN.CENTER, font=FONT_CODE)
    _text(s, "landing_page_load_time = max(0.5, 4.0 - ROAS \u00d7 0.1)",
          Inches(2), Inches(2.8), Inches(9.3), Inches(0.6),
          sz=22, color=C_ORANGE, align=PP_ALIGN.CENTER, font=FONT_CODE)
    _text(s, "creative_impact_factor = f(bounce_rate, load_time)",
          Inches(2), Inches(3.4), Inches(9.3), Inches(0.6),
          sz=22, color=C_ORANGE, align=PP_ALIGN.CENTER, font=FONT_CODE)

    _text(s, "정답(ROAS)의 역함수가 피처로 사용됨 → 모델이 정답지를 \"컨닝\"",
          Inches(1), Inches(5.0), Inches(11.3), Inches(0.6),
          sz=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _text(s, "creative_data_enricher.py에서 발견 — 5단계 Ablation Study로 추적",
          Inches(1), Inches(5.8), Inches(11.3), Inches(0.5),
          sz=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    _notes(s,
        "[나레이션]\n"
        "Phase 3의 creative_data_enricher.py에서\n"
        "bounce_rate = 65 - ROAS × 2라는 공식이 사용되고 있었습니다.\n"
        "ROAS가 높으면 bounce_rate가 낮아지는 역함수 구조 —\n"
        "모델이 이 관계를 학습하면 사실상 정답지를 보고 있는 셈입니다.")


# ── Slide 5: Waterfall ───────────────────────────────────────────────
def s05_waterfall(prs):
    s = _slide(prs, C_LIGHT)
    _accent_bar(s)
    _text(s, "R\u00b2 하락 궤적: Leakage 제거 과정",
          Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
          sz=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _fig(s, "waterfall_r2_leakage.png", y=Inches(1.0))
    _text(s, "Full Model(R\u00b2 0.79)에서 leakage 변수를 하나씩 제거 → R\u00b2 0.35로 교정",
          Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
          sz=14, color=C_DIM, align=PP_ALIGN.CENTER)
    _notes(s,
        "[나레이션 0:25-0:33]\n"
        "leakage 변수 3개를 하나씩 제거하자,\n"
        "R²는 0.79에서 0.35로 절반 이하로 떨어졌습니다.")


# ── Slide 6: SHAP ────────────────────────────────────────────────────
def s06_shap(prs):
    s = _slide(prs, C_LIGHT)
    _accent_bar(s)
    _text(s, "SHAP 기여도 대조: Leakage 포함 vs 제거",
          Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
          sz=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _fig(s, "shap_force_leakage_vs_clean.png", y=Inches(1.0))
    _text(s, "bounce_rate 독점 → 제거 후 광고주가 통제 가능한 변수들이 상위로 전환",
          Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
          sz=14, color=C_DIM, align=PP_ALIGN.CENTER)
    _notes(s,
        "[나레이션 0:33-0:41]\n"
        "SHAP 분석에서도 bounce_rate가 혼자 50%를 차지하던 게,\n"
        "제거 후에는 광고주가 통제 가능한 변수들이 상위를 차지합니다.")


# ── Slide 7: Ablation ────────────────────────────────────────────────
def s07_ablation(prs):
    s = _slide(prs, C_LIGHT)
    _accent_bar(s)
    _text(s, "Ablation Study V1→V5: 5단계 자체 감사",
          Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
          sz=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _fig(s, "ablation_guardrail.png", y=Inches(1.0))
    _text(s, "Leave-One-Out → Shapley 분해 → Robustness → 개별 피처 분해 → Data Audit",
          Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
          sz=14, color=C_DIM, align=PP_ALIGN.CENTER)
    _notes(s,
        "[나레이션 — 긴 버전용, 1분 영상에서는 SKIP 가능]\n"
        "5단계에 걸친 ablation study로 leakage를 체계적으로 추적했습니다.\n"
        "V1에서 시작해 V5 Data Generation Audit에서 정확한 공식을 발견했습니다.")


# ── Slide 8: LIVE DEMO ──────────────────────────────────────────────
def s08_demo(prs):
    s = _slide(prs, C_NAVY)

    _text(s, "LIVE DEMO",
          Inches(1), Inches(1.2), Inches(11.3), Inches(1.5),
          sz=72, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _text(s, "Streamlit AI 에이전트 실시간 시연",
          Inches(1), Inches(3.0), Inches(11.3), Inches(0.8),
          sz=28, color=C_BLUE, align=PP_ALIGN.CENTER)

    _rect(s, Inches(3.5), Inches(4.2), Inches(6.3), Inches(2.5), C_CARD)
    _multi(s, [
        ("시연 시나리오", 18, True),
        ('"핀테크 앱 광고를 미국에서 시작하려고 해요"', 16, False),
        ("→ AI 에이전트가 대화하며 전략 설계", 14, False),
        ("→ 예측 결과 차트 실시간 생성", 14, False),
        ("→ 예산 시뮬레이터 탭 전환 시연", 14, False),
    ], Inches(3.8), Inches(4.4), Inches(5.7), Inches(2.2),
       color=C_GRAY, align=PP_ALIGN.LEFT)

    _notes(s,
        "[라이브 데모 — 30~45초]\n\n"
        "1. Streamlit 앱 열기 (adstrategy-ai.streamlit.app)\n"
        "2. 사이드바 Quick Start → '핀테크 앱 광고를 미국에서 시작하려고 해요' 클릭\n"
        "3. AI 에이전트 후속 질문에 답변: '월 $5,000 예산, 3월에 시작합니다'\n"
        "4. 예측 결과 차트가 자동 생성되는 것을 보여줌\n"
        "5. 예산 시뮬레이터 탭 전환 — 슬라이더 조작하며 반응 곡선 시연\n"
        "6. (선택) Leakage 감사 탭 잠깐 보여주기\n\n"
        "TIP: 네트워크 이슈 대비 → 미리 녹화한 백업 영상 준비")


# ── Slide 9: Budget Impact ───────────────────────────────────────────
def s09_budget(prs):
    s = _slide(prs, C_LIGHT)
    _accent_bar(s)
    _text(s, "예산 재할당 효과: 상위 20% 집중 시 ROAS +170%",
          Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
          sz=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _fig(s, "budget_reallocation_impact.png", y=Inches(1.0))
    _text(s, "정직한 모델(R\u00b2\u22480.35)로도 예산 재할당만으로 의미 있는 ROAS 개선 가능",
          Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
          sz=14, color=C_DIM, align=PP_ALIGN.CENTER)
    _notes(s,
        "[나레이션 0:41-0:50]\n"
        "이 정직한 모델로 예산 재할당을 시뮬레이션하면,\n"
        "상위 20% 캠페인에 집중 시 ROAS가 170% 개선됩니다.")


# ── Slide 10: 3 Deliverables ─────────────────────────────────────────
def s10_deliverables(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)
    _text(s, "3가지 핵심 산출물",
          Inches(1), Inches(0.3), Inches(11.3), Inches(0.8),
          sz=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    cards = [
        ("1", "데이터 감사\n파이프라인",
         "합성 데이터 프로젝트에서\n재사용 가능한 5단계\nLeakage 탐지 템플릿", C_RED),
        ("2", "전략 의사결정\n도구",
         "예산 재할당 시뮬레이션\n플랫폼별 반응 곡선\nROAS +170% 개선", C_BLUE),
        ("3", "AI 에이전트",
         "GPT-4o Function Calling\nStreamlit 기반\n자연어 광고 전략 설계", C_GREEN),
    ]

    cw, ch = Inches(3.4), Inches(3.8)
    xs = [Inches(1.3), Inches(5.0), Inches(8.7)]
    ty = Inches(1.6)

    for i, (num, title, desc, clr) in enumerate(cards):
        _rect(s, xs[i], ty, cw, ch, C_CARD, border=clr, bw=Pt(2))
        _text(s, num, xs[i], ty + Inches(0.2), cw, Inches(0.6),
              sz=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
        _text(s, title, xs[i], ty + Inches(0.9), cw, Inches(0.9),
              sz=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _text(s, desc,
              xs[i] + Inches(0.3), ty + Inches(2.0),
              cw - Inches(0.6), Inches(1.5),
              sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    _notes(s,
        "[나레이션 — 긴 버전용]\n"
        "이 프로젝트의 3가지 핵심 산출물입니다.\n"
        "첫째, 합성 데이터 프로젝트에서 재사용 가능한 5단계 leakage 탐지 템플릿.\n"
        "둘째, 예산 재할당 시뮬레이션과 플랫폼별 반응 곡선.\n"
        "셋째, GPT-4o 기반 자연어 광고 전략 설계 AI 에이전트.")


# ── Slide 11: Transparency ───────────────────────────────────────────
def s11_transparency(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)
    _text(s, "데이터 투명성",
          Inches(1), Inches(0.3), Inches(11.3), Inches(0.7),
          sz=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("원본 데이터", "Kaggle 공개 광고 캠페인 1,800건"),
        ("보강 후", "10,030건 \u00d7 42피처 (4단계 파이프라인)"),
        ("합성 비율", "약 46% — 볼륨 확보 목적"),
        ("Honest R\u00b2", "0.35 (leakage 제거, 5-fold TimeSeriesSplit CV)"),
        ("모델 용도", "개별 수치 예측이 아닌, 플랫폼\u00b7캠페인 간 상대 비교"),
        ("Leakage 변수", "bounce_rate, load_time, creative_impact_factor"),
    ]
    sy = Inches(1.4)
    rh = Inches(0.78)

    for i, (label, value) in enumerate(rows):
        ry = sy + i * rh
        if i % 2 == 0:
            _rect(s, Inches(1), ry, Inches(11.3), rh, C_CARD)
        _text(s, label,
              Inches(1.4), ry + Inches(0.15), Inches(2.5), rh - Inches(0.3),
              sz=16, bold=True, color=C_BLUE, align=PP_ALIGN.LEFT)
        _text(s, value,
              Inches(4.2), ry + Inches(0.15), Inches(7.8), rh - Inches(0.3),
              sz=16, color=C_WHITE, align=PP_ALIGN.LEFT)

    _text(s, "이 프로젝트의 핵심 산출물은 예측 모델이 아니라, 데이터 감사 프레임워크입니다.",
          Inches(1), Inches(6.4), Inches(11.3), Inches(0.5),
          sz=16, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    _notes(s,
        "[나레이션 — Q&A 대비]\n"
        "데이터의 한계를 투명하게 공개합니다.\n"
        "합성 비율이 46%이고, honest R²는 0.35입니다.\n"
        "하지만 이 프로젝트의 가치는 예측 정확도가 아니라,\n"
        "높은 R²의 함정을 발견하고 교정한 감사 프레임워크에 있습니다.")


# ── Slide 12: Key Message ────────────────────────────────────────────
def s12_key_message(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)

    _text(s, "\u201c",
          Inches(1.5), Inches(0.8), Inches(2), Inches(2),
          sz=120, bold=True, color=RGBColor(0x30, 0x30, 0x40),
          align=PP_ALIGN.LEFT)
    _text(s, "높은 R\u00b2는 좋은 모델의 증거가 아니라\n감사의 출발점이다",
          Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5),
          sz=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _text(s, "R\u00b2 0.79의 하락이 실패가 아니라, 정직해진 전환점이었습니다.",
          Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.8),
          sz=22, color=C_GRAY, align=PP_ALIGN.CENTER)

    _notes(s,
        "[나레이션 0:50-1:00]\n"
        "R² 0.79의 하락이 실패가 아니라 정직해진 전환점이었습니다.\n"
        "이 프로젝트의 산출물은 예측 모델이 아니라,\n"
        "누수를 잡아내는 데이터 감사 프레임워크 그 자체입니다.")


# ── Slide 13: Thank You ──────────────────────────────────────────────
def s13_thank_you(prs):
    s = _slide(prs, C_DARK)
    _accent_bar(s)

    _text(s, "감사합니다",
          Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
          sz=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _multi(s, [
        ("Live Demo", 20, True),
        ("adstrategy-ai.streamlit.app", 18, False),
        ("", 10, False),
        ("GitHub", 20, True),
        ("github.com/Sagak-bit/AdStrategy-AI", 18, False),
    ], Inches(3.5), Inches(3.5), Inches(6.3), Inches(2.8),
       color=C_BLUE, align=PP_ALIGN.CENTER)

    _text(s, "AdStrategy AI — 데이터 기반 광고 전략 설계",
          Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
          sz=14, color=C_GRAY, align=PP_ALIGN.CENTER)

    _notes(s, "[엔딩] 감사 인사. Live Demo URL과 GitHub URL 노출.")


# ── Appendix Slides ──────────────────────────────────────────────────
def s14_appendix_leakage(prs):
    s = _slide(prs, C_LIGHT)
    _accent_bar(s, C_GRAY)
    _text(s, "[Appendix] 피처별 Leakage 위험도 대시보드",
          Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
          sz=24, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _fig(s, "leakage_risk_dashboard.png", y=Inches(1.0))
    _text(s, "각 피처를 단독으로 ROAS 예측에 사용했을 때의 R² — 비정상적으로 높으면 leakage 의심",
          Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
          sz=13, color=C_DIM, align=PP_ALIGN.CENTER)
    _notes(s, "[Q&A 대비] 피처 단독 R² 기준으로 leakage를 탐지하는 방법론 설명.")


def s15_appendix_segment(prs):
    s = _slide(prs, C_LIGHT)
    _accent_bar(s, C_GRAY)
    _text(s, "[Appendix] 세그먼트 신뢰도 히트맵",
          Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
          sz=24, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _fig(s, "segment_confidence_heatmap.png", y=Inches(1.0))
    _text(s, "플랫폼 × 산업 조합별 예측 MAE — 낮을수록 해당 조합의 예측 신뢰도가 높음",
          Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
          sz=13, color=C_DIM, align=PP_ALIGN.CENTER)
    _notes(s, "[Q&A 대비] 어떤 플랫폼×산업 조합에서 모델이 가장 신뢰할 수 있는지 보여줌.")


# ── Main ─────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    builders = [
        s01_title,
        s02_vision,
        s03_built,
        s04_hook,
        s03_pipeline,
        s04_smoking_gun,
        s05_waterfall,
        s06_shap,
        s07_ablation,
        s08_demo,
        s09_budget,
        s10_deliverables,
        s11_transparency,
        s12_key_message,
        s13_thank_you,
        s14_appendix_leakage,
        s15_appendix_segment,
    ]

    for fn in builders:
        fn(prs)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        _text(slide, f"{i}/{total}",
              Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.3),
              sz=9, color=C_GRAY, align=PP_ALIGN.RIGHT)

    output = os.path.join(ROOT, "AdStrategy_AI_발표자료_v2.pptx")
    prs.save(output)
    print(f"[OK] 발표자료 생성 완료: {output}")
    print(f"     총 {total}장 슬라이드 ({total - 2} 본편 + 2 Appendix)")
    print()
    print("슬라이드 구성:")
    names = [
        "타이틀", "★ 프로젝트 비전 (NEW)", "★ 무엇을 만들었나 (NEW)",
        "후킹 (R²=0.79)", "파이프라인", "Smoking Gun (코드)",
        "Waterfall R²", "SHAP 대조", "Ablation V1→V5", "LIVE DEMO",
        "예산 재할당", "3가지 산출물", "데이터 투명성", "핵심 메시지",
        "감사합니다", "[Appendix] Leakage 대시보드", "[Appendix] 세그먼트 히트맵",
    ]
    for i, name in enumerate(names, 1):
        marker = "  ★" if name == "LIVE DEMO" else ""
        print(f"  {i:2d}. {name}{marker}")


# =====================================================================
# 1분 전용 — Modern Dark UI (AI Dashboard 스타일)
# =====================================================================

M_BG = RGBColor(0x0A, 0x0B, 0x10)
M_CARD = RGBColor(0x13, 0x15, 0x1C)
M_CARD_BRD = RGBColor(0x2A, 0x2D, 0x39)
M_PURPLE = RGBColor(0x6C, 0x5C, 0xE7)
M_TEAL = RGBColor(0x00, 0xD2, 0xD3)
M_CORAL = RGBColor(0xFF, 0x76, 0x75)
M_MINT = RGBColor(0x55, 0xEF, 0xC4)
M_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
M_SUB = RGBColor(0xA0, 0xA5, 0xB5)
M_DIM = RGBColor(0x5A, 0x5F, 0x73)
M_CODE_BG = RGBColor(0x05, 0x05, 0x07)

FT = "Pretendard"
FT_B = "Pretendard"


def _mcard(s, x, y, w, h, fill=M_CARD, border=M_CARD_BRD, bw_=Pt(1)):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = bw_
    else:
        shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.05
    except Exception:
        pass
    return shape


def _accent(s, x, y, w, h, color):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _mfooter(s, num, total):
    _text(s, f"{num} / {total}",
          SW - Inches(1.4), SH - Inches(0.4), Inches(1), Inches(0.3),
          sz=9, color=M_DIM, align=PP_ALIGN.RIGHT, font=FT)


def c01_title_vision(prs):
    """Slide 1: 타이틀 + 가설"""
    s = _slide(prs, M_BG)

    _text(s, "AdStrategy AI",
          Inches(1.5), Inches(1.5), Inches(8), Inches(1),
          sz=64, bold=True, color=M_WHITE, font=FT_B)

    _accent(s, Inches(1.5), Inches(2.8), Pt(4), Inches(0.35), M_PURPLE)
    _text(s, "HYPOTHESIS",
          Inches(1.75), Inches(2.75), Inches(3), Inches(0.4),
          sz=11, bold=True, color=M_PURPLE, font=FT)

    _mcard(s, Inches(1.5), Inches(3.4), Inches(10.3), Inches(2.2))
    _accent(s, Inches(1.5), Inches(3.4), Inches(10.3), Pt(3), M_PURPLE)
    _text(s, "광고 예산 배분을\n데이터로 최적화할 수 있는가?",
          Inches(1.8), Inches(3.7), Inches(9.7), Inches(1.5),
          sz=34, bold=True, color=M_WHITE, font=FT_B)
    _text(s, "Kaggle 1,800건 \u2192 4단계 보강 \u2192 ML 예측 \u2192 Ablation 감사 \u2192 AI 에이전트까지",
          Inches(1.8), Inches(5.0), Inches(9.7), Inches(0.4),
          sz=14, color=M_SUB, font=FT)

    # 하단 키워드 3개
    items = [
        ("DATA", "10,030건 \u00d7 42피처", M_PURPLE),
        ("MODEL", "XGBoost 앙상블", M_TEAL),
        ("AGENT", "GPT-4o Function Calling", M_MINT),
    ]
    cw = Inches(3.2)
    xs = [Inches(1.5), Inches(5.1), Inches(8.7)]
    for i, (tag, desc, clr) in enumerate(items):
        _mcard(s, xs[i], Inches(5.8), cw, Inches(1.2))
        _accent(s, xs[i], Inches(5.8), cw, Pt(3), clr)
        _text(s, tag, xs[i] + Inches(0.25), Inches(5.95), cw, Inches(0.3),
              sz=10, bold=True, color=clr, font=FT)
        _text(s, desc, xs[i] + Inches(0.25), Inches(6.25), cw, Inches(0.4),
              sz=14, bold=True, color=M_WHITE, font=FT_B)

    _notes(s,
        "[0:00-0:07] 7초\n"
        "광고 예산 배분을 데이터로 최적화할 수 있는가.\n"
        "공개 데이터 만 건을 보강하고, 모델과 에이전트까지 연결한 프로젝트입니다.")


def c02_built(prs):
    """Slide 2: 엔드투엔드 파이프라인 다이어그램"""
    s = _slide(prs, M_BG)

    _text(s, "End-to-End Pipeline",
          Inches(1.0), Inches(0.4), Inches(10), Inches(0.7),
          sz=32, bold=True, color=M_WHITE, font=FT_B)
    _text(s, "데이터 수집부터 광고주 의사결정까지 하나의 파이프라인으로 연결",
          Inches(1.0), Inches(1.05), Inches(10), Inches(0.4),
          sz=14, color=M_SUB, font=FT)

    # ── 상단: 6단계 파이프라인 (가로 플로우) ──
    stages = [
        ("01", "데이터 수집", "Kaggle\n1,800건", M_PURPLE),
        ("02", "4단계 보강", "API·합성\n10,030건", M_PURPLE),
        ("03", "ML 모델링", "XGBoost\nROAS 예측", M_TEAL),
        ("04", "Ablation 감사", "5단계 검증\nLeakage 탐지", M_CORAL),
        ("05", "AI 에이전트", "GPT-4o\nFunction Call", M_MINT),
        ("06", "광고주 활용", "Streamlit\n전략 설계", RGBColor(0xFD, 0xCB, 0x6E)),
    ]
    bw_, bh_ = Inches(1.7), Inches(2.5)
    gap = Inches(0.25)
    start_x = Inches(0.35)
    ty = Inches(1.8)

    for i, (num, title, desc, clr) in enumerate(stages):
        x = start_x + i * (bw_ + gap)
        border = M_CORAL if i == 3 else M_CARD_BRD
        bw_val = Pt(2) if i == 3 else Pt(1)
        _mcard(s, x, ty, bw_, bh_, border=border, bw_=bw_val)
        _accent(s, x, ty, bw_, Pt(3), clr)

        _text(s, num, x + Inches(0.12), ty + Inches(0.2), Inches(0.5), Inches(0.3),
              sz=11, bold=True, color=clr, font=FT)
        _text(s, title, x + Inches(0.12), ty + Inches(0.55), bw_ - Inches(0.2), Inches(0.5),
              sz=13, bold=True, color=M_WHITE, font=FT_B)
        _text(s, desc, x + Inches(0.12), ty + Inches(1.1), bw_ - Inches(0.2), Inches(1.0),
              sz=11, color=M_SUB, font=FT)

        if i < 5:
            ax = x + bw_ + Inches(0.03)
            _text(s, "\u2192", ax, ty + Inches(0.8), gap - Inches(0.06), Inches(0.4),
                  sz=18, color=M_DIM, align=PP_ALIGN.CENTER, font=FT)

    # Stage 04에 경고 표시
    x4 = start_x + 3 * (bw_ + gap)
    _text(s, "\u26a0", x4 + bw_ - Inches(0.35), ty + Inches(0.12), Inches(0.3), Inches(0.3),
          sz=14, color=M_CORAL, align=PP_ALIGN.CENTER, font=FT)

    # ── 하단: "그런데 —" ──
    _mcard(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.6))
    _accent(s, Inches(1.0), Inches(4.8), Inches(11.3), Pt(3), M_CORAL)
    _text(s, "여기까지는 순조로웠습니다.   그런데 —",
          Inches(1.3), Inches(5.0), Inches(10.7), Inches(0.5),
          sz=24, bold=True, color=M_CORAL, font=FT_B)
    _text(s, "04단계 Ablation 감사에서 심각한 문제를 발견합니다",
          Inches(1.3), Inches(5.5), Inches(10.7), Inches(0.4),
          sz=14, color=M_SUB, font=FT)

    _mfooter(s, 2, 7)
    _notes(s,
        "[0:07-0:15] 8초\n"
        "공개 데이터를 4단계로 만 건까지 보강하고,\n"
        "모델, 감사, 에이전트까지 하나의 파이프라인으로 연결했습니다.\n"
        "그런데 —")


def c03_problem(prs):
    s = _slide(prs, M_BG)

    # 좌측: R² = 0.79
    _mcard(s, Inches(0.5), Inches(0.5), Inches(5.7), Inches(6.2))
    _text(s, "R\u00b2 = 0.79",
          Inches(0.5), Inches(1.5), Inches(5.7), Inches(1.5),
          sz=68, bold=True, color=M_CORAL, font=FT_B, align=PP_ALIGN.CENTER)
    _text(s, "마케팅 데이터에서\n비정상적으로 높은 수치",
          Inches(0.5), Inches(3.5), Inches(5.7), Inches(1.0),
          sz=20, color=M_WHITE, font=FT, align=PP_ALIGN.CENTER)
    _text(s, "Meta Robyn MMM 벤치마크: R\u00b2 0.50\u007e0.70",
          Inches(0.5), Inches(5.0), Inches(5.7), Inches(0.5),
          sz=12, color=M_DIM, font=FT, align=PP_ALIGN.CENTER)

    # 우측: 코드 에디터 스타일
    _mcard(s, Inches(6.8), Inches(0.5), Inches(6.1), Inches(6.2),
           fill=M_CODE_BG, border=M_CORAL, bw_=Pt(2))
    # 에디터 상단 버튼
    for j, c in enumerate([M_CORAL, RGBColor(0xFD, 0xCB, 0x6E), M_MINT]):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(7.05) + j * Inches(0.25), Inches(0.75),
                                 Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()

    _text(s, "Phase 3에서 발견된 Data Leakage",
          Inches(6.8), Inches(1.1), Inches(6.1), Inches(0.5),
          sz=13, bold=True, color=M_CORAL, font=FT, align=PP_ALIGN.CENTER)
    _text(s, "bounce_rate = 65 - ROAS * 2",
          Inches(7.0), Inches(2.3), Inches(5.7), Inches(0.6),
          sz=26, bold=True, color=M_CORAL, font=FONT_CODE, align=PP_ALIGN.LEFT)
    _text(s, "landing_page_load_time = max(0.5, 4.0 - ROAS * 0.1)\n"
          "creative_impact_factor = f(bounce_rate, load_time)",
          Inches(7.0), Inches(3.2), Inches(5.7), Inches(1.2),
          sz=13, color=M_SUB, font=FONT_CODE, align=PP_ALIGN.LEFT)
    _text(s, "정답(ROAS)의 역함수가 피처로 사용됨",
          Inches(6.8), Inches(4.8), Inches(6.1), Inches(0.5),
          sz=16, color=M_WHITE, font=FT, align=PP_ALIGN.CENTER)
    _text(s, "Target Leakage — 정답지를 보고 시험 치는 셈",
          Inches(6.8), Inches(5.5), Inches(6.1), Inches(0.5),
          sz=13, color=M_DIM, font=FT, align=PP_ALIGN.CENTER)

    _mfooter(s, 3, 7)
    _notes(s,
        "[0:15-0:25] 10초\n"
        "예측 모델의 R²가 0.79로 비정상적으로 높았습니다.\n"
        "체계적 감사로 원인을 추적해보니,\n"
        "정답을 암시하는 데이터 누수 변수가 학습에 섞여 있었습니다.")


def c04_evidence(prs):
    s = _slide(prs, M_BG)

    _text(s, "Leakage 제거  \u2192  R\u00b2 0.79에서 0.35로",
          Inches(1.0), Inches(0.3), Inches(11.3), Inches(0.6),
          sz=24, bold=True, color=M_WHITE, font=FT_B, align=PP_ALIGN.LEFT)
    _accent(s, Inches(1.0), Inches(0.95), Inches(3), Pt(3), M_PURPLE)

    _mcard(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.5))
    _fig(s, "waterfall_r2_leakage.png", y=Inches(1.3), max_h=Inches(5.2),
         max_w=Inches(11))

    _text(s, "비로소 정답지를 빼고 푼, 모델의 '진짜 실력'을 확인",
          Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.4),
          sz=13, color=M_DIM, font=FT, align=PP_ALIGN.LEFT)

    _mfooter(s, 4, 7)
    _notes(s,
        "[0:25-0:33] 8초\n"
        "해당 변수들을 모두 제거하자 R²는 0.35로 떨어졌습니다.\n"
        "비로소 정답지를 빼고 푼, 모델의 '진짜 실력'을 확인한 겁니다.")


def c05_demo(prs):
    """Slide 5: 분석→에이전트 연결 파이프라인 + 데모 전환"""
    s = _slide(prs, M_BG)

    _text(s, "분석 결과를 어떻게 활용하는가",
          Inches(1.0), Inches(0.3), Inches(10), Inches(0.6),
          sz=26, bold=True, color=M_WHITE, font=FT_B)
    _text(s, "정직한 모델 \u2192 AI 에이전트 \u2192 광고주 의사결정",
          Inches(1.0), Inches(0.85), Inches(10), Inches(0.4),
          sz=13, color=M_SUB, font=FT)

    # ── 상단: 분석→에이전트 미니 파이프라인 (5단계) ──
    flow = [
        ("Honest Model", "R\u00b2=0.35", M_TEAL),
        ("GPT-4o", "Function\nCalling", M_PURPLE),
        ("4 Tools", "\uc608\uce21 \u00b7 \ube44\uad50\n\ubca4\uce58\ub9c8\ud06c \u00b7 \ud2b8\ub80c\ub4dc", M_PURPLE),
        ("Streamlit", "5\ud0ed\n\ub300\uc2dc\ubcf4\ub4dc", M_MINT),
        ("\uad11\uace0\uc8fc", "\uc790\uc5f0\uc5b4 \ub300\ud654\ub85c\n\uc804\ub7b5 \uc124\uacc4", RGBColor(0xFD, 0xCB, 0x6E)),
    ]
    fw, fh = Inches(2.0), Inches(1.6)
    fgap = Inches(0.3)
    fstart = Inches(0.6)
    fy = Inches(1.5)

    for i, (title, desc, clr) in enumerate(flow):
        fx = fstart + i * (fw + fgap)
        _mcard(s, fx, fy, fw, fh)
        _accent(s, fx, fy, fw, Pt(3), clr)
        _text(s, title, fx + Inches(0.1), fy + Inches(0.2), fw - Inches(0.2), Inches(0.4),
              sz=13, bold=True, color=clr, font=FT)
        _text(s, desc, fx + Inches(0.1), fy + Inches(0.65), fw - Inches(0.2), Inches(0.7),
              sz=10, color=M_SUB, font=FT)
        if i < 4:
            ax = fx + fw + Inches(0.02)
            _text(s, "\u2192", ax, fy + Inches(0.4), fgap - Inches(0.04), Inches(0.4),
                  sz=16, color=M_DIM, align=PP_ALIGN.CENTER, font=FT)

    # ── 중단: DEMO 전환 영역 ──
    _mcard(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(2.5),
           fill=RGBColor(0x10, 0x12, 0x1A), border=M_PURPLE, bw_=Pt(2))
    _text(s, "LIVE DEMO",
          Inches(1.0), Inches(3.7), Inches(11.3), Inches(1.0),
          sz=52, bold=True, color=M_PURPLE, font=FT_B, align=PP_ALIGN.CENTER)
    _text(s, "adstrategy-ai.streamlit.app",
          Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.5),
          sz=18, color=M_TEAL, font=FT, align=PP_ALIGN.CENTER)

    # ── 하단: 편집 가이드 ──
    _mcard(s, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.9))
    _text(s, "\ud3b8\uc9d1: \ud0c0\uc774\ud551 10x \u2192 \ub300\uae30 \ucef7 \u2192 \uc751\ub2f5 3x \u2192 \ud0ed \uc804\ud658 5x   |   \uc774 \uc2ac\ub77c\uc774\ub4dc 2\ucd08 \ud6c4 \ub370\ubaa8 \uc601\uc0c1 \uc804\ud658",
          Inches(2.7), Inches(6.4), Inches(7.9), Inches(0.6),
          sz=10, color=M_DIM, font=FT)

    _mfooter(s, 5, 7)
    _notes(s,
        "[0:33-0:47] 데모 14초\n\n"
        "이 분석 결과를 AI 에이전트와 연결해,\n"
        "광고주가 대화만으로 데이터 기반 전략을 받을 수 있게 했습니다.\n\n"
        "이 슬라이드 2초 노출 후 Streamlit 녹화 영상 전환.\n"
        "① 타이핑 10x → ② 대기 컷 → ③ 응답 3x → ④ 차트 원속\n"
        "⑤ 탭전환 5x → ⑥ 결과 정지 원속  →  총 ~12초")


def c06_impact(prs):
    s = _slide(prs, M_BG)

    _text(s, "+170%",
          Inches(1), Inches(0.5), Inches(11.3), Inches(2.2),
          sz=100, bold=True, color=M_MINT, font=FT_B, align=PP_ALIGN.CENTER)
    _text(s, "정직해진 모델로 상위 캠페인에 예산 집중 시  ROAS 개선",
          Inches(1), Inches(2.8), Inches(11.3), Inches(0.8),
          sz=20, color=M_SUB, font=FT, align=PP_ALIGN.CENTER)

    _mcard(s, Inches(1.5), Inches(4.2), Inches(4.6), Inches(2.2))
    _accent(s, Inches(1.5), Inches(4.2), Inches(4.6), Pt(3), M_TEAL)
    _text(s, "ACCURACY", Inches(1.8), Inches(4.5), Inches(4), Inches(0.3),
          sz=10, bold=True, color=M_TEAL, font=FT)
    _text(s, "R\u00b2 0.35여도\n방향은 맞습니다",
          Inches(1.8), Inches(4.9), Inches(4), Inches(1.0),
          sz=20, bold=True, color=M_WHITE, font=FT_B)

    _mcard(s, Inches(7.2), Inches(4.2), Inches(4.6), Inches(2.2))
    _accent(s, Inches(7.2), Inches(4.2), Inches(4.6), Pt(3), M_MINT)
    _text(s, "VALUE", Inches(7.5), Inches(4.5), Inches(4), Inches(0.3),
          sz=10, bold=True, color=M_MINT, font=FT)
    _text(s, "어떤 캠페인이 더 나은지\n가리키는 데는 충분",
          Inches(7.5), Inches(4.9), Inches(4), Inches(1.0),
          sz=20, bold=True, color=M_WHITE, font=FT_B)

    _mfooter(s, 6, 7)
    _notes(s,
        "[0:47-0:53] 6초\n"
        "이 정직해진 모델로 상위 캠페인에 예산을 집중 시뮬레이션한 결과,\n"
        "ROAS를 170%까지 개선할 수 있었습니다.")


def c07_end(prs):
    s = _slide(prs, M_BG)
    GITHUB_LOGO = os.path.join(ROOT, "githublogo.png")

    # ── 좌측: 인용문 + 메시지 (슬라이드의 60%) ──
    _accent(s, Inches(1.2), Inches(0.8), Pt(4), Inches(2.8), M_PURPLE)

    _text(s, "높은 R\u00b2는\n좋은 모델의 증거가 아니라\n감사의 출발점이다",
          Inches(1.5), Inches(0.7), Inches(6.5), Inches(2.8),
          sz=38, bold=True, color=M_WHITE, font=FT_B)

    _text(s, "AI의 오류를 잡아내는 건 결국 사람의 몫입니다.\n"
          "정직한 데이터가 진짜 가치를 만듭니다.",
          Inches(1.5), Inches(3.8), Inches(6.5), Inches(1.0),
          sz=18, color=M_SUB, font=FT)

    _text(s, "Thank You",
          Inches(1.5), Inches(5.5), Inches(6.5), Inches(1.0),
          sz=52, bold=True, color=RGBColor(0x2A, 0x2D, 0x39), font="Freesentation")

    # ── 우측: 프로젝트 링크 카드 (슬라이드의 35%) ──
    rx = Inches(8.5)
    rw = Inches(4.3)

    # Live Demo 카드
    _mcard(s, rx, Inches(1.0), rw, Inches(2.2))
    _accent(s, rx, Inches(1.0), rw, Pt(3), M_TEAL)
    _text(s, "LIVE DEMO", rx + Inches(0.3), Inches(1.3), rw, Inches(0.3),
          sz=11, bold=True, color=M_TEAL, font=FT)
    _text(s, "adstrategy-ai\n.streamlit.app",
          rx + Inches(0.3), Inches(1.7), rw - Inches(0.5), Inches(0.9),
          sz=20, bold=True, color=M_WHITE, font=FT_B)
    _text(s, "5탭 대시보드 · AI 에이전트 · 예산 시뮬레이터",
          rx + Inches(0.3), Inches(2.6), rw - Inches(0.5), Inches(0.4),
          sz=10, color=M_DIM, font=FT)

    # GitHub 카드
    _mcard(s, rx, Inches(3.5), rw, Inches(2.2))
    _accent(s, rx, Inches(3.5), rw, Pt(3), M_PURPLE)
    _text(s, "SOURCE CODE", rx + Inches(0.3), Inches(3.8), rw, Inches(0.3),
          sz=11, bold=True, color=M_PURPLE, font=FT)

    if os.path.exists(GITHUB_LOGO):
        try:
            s.shapes.add_picture(GITHUB_LOGO,
                                 rx + Inches(0.3), Inches(4.2),
                                 Inches(0.4), Inches(0.4))
        except Exception:
            pass
    _text(s, "Sagak-bit /\nAdStrategy-AI",
          rx + Inches(0.9), Inches(4.15), rw - Inches(1.2), Inches(0.8),
          sz=18, bold=True, color=M_WHITE, font=FT_B)
    _text(s, "파이프라인 · ML 모델 · 에이전트 · 분석 전체 코드",
          rx + Inches(0.3), Inches(5.1), rw - Inches(0.5), Inches(0.4),
          sz=10, color=M_DIM, font=FT)

    # 하단 크레딧
    _text(s, "AdStrategy AI  —  데이터 기반 광고 전략 설계",
          Inches(1.5), Inches(6.8), Inches(11), Inches(0.4),
          sz=11, color=M_DIM, font=FT, align=PP_ALIGN.CENTER)

    _mfooter(s, 7, 7)
    _notes(s,
        "[0:53-1:00] 7초\n"
        "AI의 오류를 잡아내는 건 결국 사람의 몫입니다.\n"
        "정직한 데이터가 진짜 가치를 만듭니다.\n"
        "감사합니다.")


def main_1min():
    """1분 전용 슬라이드 생성 (7장)"""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    builders = [
        c01_title_vision,
        c02_built,
        c03_problem,
        c04_evidence,
        c05_demo,
        c06_impact,
        c07_end,
    ]

    for fn in builders:
        fn(prs)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        _text(slide, f"{i}/{total}",
              Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.3),
              sz=9, color=C_GRAY, align=PP_ALIGN.RIGHT)

    output = os.path.join(ROOT, "AdStrategy_AI_1min_DA.pptx")
    prs.save(output)
    print(f"[OK] 1분 전용 발표자료: {output}")
    print(f"     총 {total}장 슬라이드")
    print()
    names = [
        "타이틀+비전", "무엇을 만들었나", "문제+코드 (좌우분할)",
        "증거 (Waterfall)", "DEMO 전환", "+170% 임팩트",
        "핵심 메시지+감사",
    ]
    print("타임라인:")
    times = ["0:00-0:07", "0:07-0:15", "0:15-0:25",
             "0:25-0:33", "0:33-0:47 (영상)", "0:47-0:53", "0:53-1:00"]
    for i, (name, t) in enumerate(zip(names, times), 1):
        print(f"  {i}. [{t}] {name}")


if __name__ == "__main__":
    if "--1min" in sys.argv:
        main_1min()
    elif "--all" in sys.argv:
        main()
        main_1min()
    else:
        main()
        main_1min()
